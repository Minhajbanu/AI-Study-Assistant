from pypdf import PdfReader
from pptx import Presentation
from PIL import Image
import pytesseract


def extract_text_from_pdf(file):

    reader = PdfReader(file)
    text = ""

    for page in reader.pages:
        text += page.extract_text()

    return text


def extract_text_from_ppt(file):

    prs = Presentation(file)
    text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text


def extract_text_from_image(file):

    image = Image.open(file)
    text = pytesseract.image_to_string(image)

    return text


def extract_text(file):

    file_type = file.name.split(".")[-1].lower()

    if file_type == "pdf":
        return extract_text_from_pdf(file)

    elif file_type == "pptx":
        return extract_text_from_ppt(file)

    elif file_type in ["png", "jpg", "jpeg"]:
        return extract_text_from_image(file)

    else:
        return ""